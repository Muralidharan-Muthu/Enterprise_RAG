"""
Generate a professional PowerPoint presentation for Multi-Store RAG Chatbot System Architecture.
Uses python-pptx with custom styling for a polished, enterprise-grade deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Color Palette (Light Mode) ─────────────────────────────────────────────
DARK_BG = RGBColor(0x1A, 0x1D, 0x29)        # Dark navy (for header bars)
CARD_BG = RGBColor(0xF5, 0xF6, 0xFA)        # Light card background
ACCENT = RGBColor(0x4A, 0x4B, 0xC9)         # Indigo/violet accent (deeper)
ACCENT_LIGHT = RGBColor(0x4A, 0x4B, 0xC9)   # Accent for text on light bg
WHITE = RGBColor(0x1A, 0x1D, 0x29)          # Primary text (dark navy)
LIGHT_GRAY = RGBColor(0x56, 0x5B, 0x6E)     # Secondary text
MEDIUM_GRAY = RGBColor(0x8A, 0x8D, 0x9F)    # Tertiary text
GREEN = RGBColor(0x15, 0x7A, 0x50)          # Success green (darker)
ORANGE = RGBColor(0xB4, 0x75, 0x0A)         # Warning/highlight (darker)
TEAL = RGBColor(0x1A, 0x7F, 0x7A)           # Teal accent (darker)
RED_SOFT = RGBColor(0xC0, 0x39, 0x2F)       # Soft red (darker)
SLIDE_BG = RGBColor(0xFF, 0xFF, 0xFF)       # White slide background
TRUE_WHITE = RGBColor(0xFF, 0xFF, 0xFF)     # Actual white for badges/pills


def set_slide_bg(slide, color=SLIDE_BG):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_with_fill(slide, left, top, width, height, color, corner_radius=None):
    """Add a rounded rectangle shape with fill."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if corner_radius is not None:
        # Set corner radius via XML
        sp = shape._element
        prstGeom = sp.find(qn('a:prstGeom'), sp.nsmap) if hasattr(sp, 'nsmap') else None
    return shape


def add_accent_bar(slide, top, width=Inches(1.2)):
    """Add a thin accent bar as a visual divider."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(tf, items, font_size=13, color=LIGHT_GRAY, bold_prefix=True):
    """Add bullet points to a text frame."""
    for item in items:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        p.level = 0

        if isinstance(item, tuple):
            # (bold_part, rest)
            run1 = p.add_run()
            run1.text = item[0]
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = WHITE
            run1.font.bold = True
            run1.font.name = "Calibri"

            run2 = p.add_run()
            run2.text = item[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.bold = False
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = f"  •  {item}"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"


def add_card(slide, left, top, width, height, title, bullets, title_color=ACCENT_LIGHT):
    """Add a card-style content block."""
    # Card background
    card = add_shape_with_fill(slide, left, top, width, height, CARD_BG)
    card.line.color.rgb = RGBColor(0xE2, 0xE4, 0xEC)
    card.line.width = Pt(1)

    # Title
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Calibri"

    # Bullets
    txBox2 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for b in bullets:
        p = tf2.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"▸ {b}"
        run.font.size = Pt(11)
        run.font.color.rgb = LIGHT_GRAY
        run.font.name = "Calibri"

    return card


def create_section_header(slide, number, title, subtitle=""):
    """Create a consistent section header with number badge."""
    set_slide_bg(slide)

    # Section number badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.35), Inches(0.55), Inches(0.45)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    badge_tf = badge.text_frame
    badge_tf.paragraphs[0].text = number
    badge_tf.paragraphs[0].font.size = Pt(16)
    badge_tf.paragraphs[0].font.color.rgb = TRUE_WHITE
    badge_tf.paragraphs[0].font.bold = True
    badge_tf.paragraphs[0].font.name = "Consolas"
    badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    add_text_box(slide, Inches(1.35), Inches(0.3), Inches(8), Inches(0.55),
                 title, font_size=26, color=WHITE, bold=True)

    # Accent bar
    add_accent_bar(slide, Inches(0.95), width=Inches(2))

    # Subtitle
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.1), Inches(8.4), Inches(0.6),
                     subtitle, font_size=13, color=LIGHT_GRAY)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title Slide
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)

    # Large accent shape (decorative)
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.8))
    deco.fill.solid()
    deco.fill.fore_color.rgb = RGBColor(0x1A, 0x1D, 0x29)
    deco.line.fill.background()

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.8), Inches(13.333), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Company branding
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                 "DECISION MINDS", font_size=14, color=RGBColor(0xB3, 0xB4, 0xFB), bold=True,
                 font_name="Calibri")

    # Main title
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(1.0),
                 "Multi-Store RAG Chatbot", font_size=48, color=TRUE_WHITE, bold=True)

    # Subtitle
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(10), Inches(0.6),
                 "Enterprise Multi-Store Retrieval-Augmented Generation System",
                 font_size=20, color=RGBColor(0xA3, 0xA6, 0xBD))

    # Description below accent line
    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(11), Inches(1.0),
                 "An intelligent document processing and retrieval system that parses, classifies, and routes "
                 "documents into specialised stores. Retrieval blends vector search, keyword search, structured "
                 "table queries, and a Neo4j knowledge graph before an LLM synthesises cited answers.",
                 font_size=14, color=LIGHT_GRAY)

    # Tech badges row
    badges = ["FastAPI", "Celery + Redis", "Supabase / pgvector", "Neo4j GraphRAG",
              "Next.js 14", "BAAI/bge-large-en-v1.5", "bge-reranker", "Gemma 4 (CDAC)"]
    x_pos = Inches(0.8)
    for badge_text in badges:
        w = Inches(len(badge_text) * 0.11 + 0.4)
        badge_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(4.8), w, Inches(0.38)
        )
        badge_shape.fill.solid()
        badge_shape.fill.fore_color.rgb = RGBColor(0xED, 0xEC, 0xFE)
        badge_shape.line.color.rgb = RGBColor(0xC5, 0xC5, 0xEE)
        badge_shape.line.width = Pt(1)
        tf = badge_shape.text_frame
        tf.paragraphs[0].text = badge_text
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = ACCENT
        tf.paragraphs[0].font.name = "Consolas"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        x_pos += w + Inches(0.15)

    # Date
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(5), Inches(0.4),
                 "System Architecture Reference  ·  July 2026", font_size=11, color=MEDIUM_GRAY)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — Overview
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "01", "Overview",
                          "Multi-Store RAG Chatbot replaces flat single-store RAG with a type-aware pipeline — every document "
                          "is parsed, classified, chunked with a type-appropriate strategy, embedded, and "
                          "stored in the Postgres store that best fits its shape.")

    add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.8),
             "Why Multi-Store?", [
                 "Different documents have different shapes — policy PDFs,",
                 "10-K financial reports, legal contracts, research papers",
                 "Flattening into generic chunks loses structure",
                 "Table rows, clause risk levels, citation metadata",
                 "are cheap to preserve and improve retrieval precision",
                 "Five specialised stores: vector, table, clause, document, image"
             ])

    add_card(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(2.8),
             "Design Principle: Fail Open", [
                 "Every smart stage has a rule-based fallback",
                 "Gemma classification → keyword-scoring rules",
                 "Semantic routing → confidence threshold guard",
                 "GraphRAG → best-effort, non-fatal",
                 "Agentic RAVEN/SPYDER → classic retrieval",
                 "Pipeline degrades gracefully — never blocks"
             ])

    # Five document types
    doc_types = [
        ("Policy", ACCENT_LIGHT), ("Financial", ORANGE),
        ("Legal", RED_SOFT), ("Research", GREEN), ("Entity", TEAL)
    ]
    x = Inches(0.8)
    for dtype, color in doc_types:
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.3), Inches(2.0), Inches(0.4))
        chip.fill.solid()
        chip.fill.fore_color.rgb = CARD_BG
        chip.line.color.rgb = color
        chip.line.width = Pt(1.5)
        tf = chip.text_frame
        tf.paragraphs[0].text = dtype
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = "Calibri"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        x += Inches(2.3)

    add_text_box(slide, Inches(0.8), Inches(5.85), Inches(11), Inches(0.4),
                 "▲  Five document types drive routing, chunking, and storage decisions across the entire pipeline",
                 font_size=11, color=MEDIUM_GRAY)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — Tech Stack
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "02", "Tech Stack",
                          "Full-stack enterprise architecture with local ML models, managed databases, and modern frontend.")

    # Backend card
    add_card(slide, Inches(0.6), Inches(2.0), Inches(3.0), Inches(3.2),
             "⚙  Backend", [
                 "FastAPI — REST API, SSE streaming",
                 "Celery — async ingestion pipeline",
                 "Redis — broker, result backend",
                 "Docling — OCR + table-structure parsing",
                 "PyMuPDF — page pre-scan, fallback"
             ], title_color=ACCENT_LIGHT)

    # Data layer card
    add_card(slide, Inches(3.85), Inches(2.0), Inches(3.0), Inches(3.2),
             "🗄  Data Layer", [
                 "Supabase Postgres — 5 chunk stores",
                 "pgvector — HNSW ANN, cosine distance",
                 "Postgres FTS — tsvector / GIN",
                 "Supabase Storage — files, images",
                 "Neo4j (Aura) — knowledge graph"
             ], title_color=GREEN)

    # ML Models card
    add_card(slide, Inches(7.1), Inches(2.0), Inches(3.0), Inches(3.2),
             "🧠  ML Models", [
                 "BAAI/bge-large-en-v1.5 — 1024-dim",
                 "ms-marco-MiniLM-L-6-v2 — reranker",
                 "EasyOCR — image OCR",
                 "Gemma 4 27B-it (CDAC) — LLM",
                 "  classify, extract, judge, synthesise"
             ], title_color=ORANGE)

    # Frontend card
    add_card(slide, Inches(10.35), Inches(2.0), Inches(2.5), Inches(3.2),
             "🖥  Frontend", [
                 "Next.js 14 — App Router",
                 "TanStack Query — polling",
                 "Tailwind CSS",
                 "API proxy → FastAPI"
             ], title_color=TEAL)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — High-Level Architecture
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "03", "High-Level Architecture",
                          "Service topology — what talks to what. Full end-to-end system view.")

    # Architecture tiers
    tiers = [
        ("Browser", "Chat UI  ·  Upload UI  ·  Document Browser", Inches(2.0), ACCENT_LIGHT),
        ("Next.js 14 Frontend (port 3000)", "App Router + API Proxy Routes", Inches(3.0), TEAL),
        ("FastAPI Backend (port 8000)", "documents · ingest · query · graph · chats · health", Inches(4.0), ACCENT_LIGHT),
    ]

    for title, sub, top, color in tiers:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(4.0), top, Inches(5.333), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(13)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = sub
        run2.font.size = Pt(9)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.name = "Calibri"

    # Arrow connectors
    for top in [Inches(2.7), Inches(3.7)]:
        arrow = add_text_box(slide, Inches(6.2), top, Inches(1), Inches(0.3),
                             "▼", font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(5.5), Inches(4.7), Inches(2.5), Inches(0.3),
                 "▼  dispatches / queries", font_size=10, color=MEDIUM_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # Bottom service tier
    services = [
        ("Celery Worker", "ingestion, community", ACCENT_LIGHT),
        ("Redis", "broker · backend", RED_SOFT),
        ("Supabase PG", "5 stores + registry", GREEN),
        ("Supabase Storage", "files, images", ORANGE),
        ("Neo4j (Aura)", "graph store", TEAL),
        ("Gemma 4 (CDAC)", "LLM endpoint", ACCENT_LIGHT),
    ]
    x = Inches(0.6)
    for svc, sub, color in services:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.2), Inches(1.95), Inches(0.8))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = svc
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = sub
        run2.font.size = Pt(8)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.name = "Calibri"
        x += Inches(2.1)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — Ingestion Pipeline (Overview)
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "04", "Ingestion Pipeline",
                          "Triggered by POST /api/v1/ingest/upload or /pipeline. Runs as a Celery task with live progress tracking.")

    # Pipeline stages as connected boxes
    stages = [
        ("0", "Cleanup", "Idempotent wipe"),
        ("1", "Parse (Docling)", "PDF / multi-format"),
        ("2", "Route", "Gemma classify → fallback"),
        ("3", "Chunk", "Type-aware splitting"),
        ("4", "Embed", "BGE 1024-dim vectors"),
        ("5", "Store", "Route to target store"),
        ("6", "Graph Stage", "Neo4j entities & rels"),
        ("7", "Completion", "Lineage checks"),
    ]

    x = Inches(0.4)
    y = Inches(2.2)
    for i, (num, title, sub) in enumerate(stages):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.4), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT if i < 7 else GREEN
        box.line.width = Pt(1.5)

        # Number badge
        nbadge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.5), y - Inches(0.15), Inches(0.35), Inches(0.35))
        nbadge.fill.solid()
        nbadge.fill.fore_color.rgb = ACCENT if i < 7 else GREEN
        nbadge.line.fill.background()
        ntf = nbadge.text_frame
        ntf.paragraphs[0].text = num
        ntf.paragraphs[0].font.size = Pt(10)
        ntf.paragraphs[0].font.color.rgb = TRUE_WHITE
        ntf.paragraphs[0].font.bold = True
        ntf.paragraphs[0].font.name = "Consolas"
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE

        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = sub
        run2.font.size = Pt(9)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.name = "Calibri"

        if i < len(stages) - 1:
            arrow = add_text_box(slide, x + Inches(1.4), y + Inches(0.3),
                                 Inches(0.25), Inches(0.3), "→", font_size=16,
                                 color=ACCENT, alignment=PP_ALIGN.CENTER)
        x += Inches(1.6)

    # Details below
    add_card(slide, Inches(0.6), Inches(3.8), Inches(5.8), Inches(3.0),
             "Parsing & Routing", [
                 "PDF: page-chunked Docling → fallback whole-doc → PyMuPDF",
                 "Non-PDF: docx, pptx, xlsx, html, md via Docling",
                 "Multi-page tables stitched; merged-cell spans detected",
                 "Gemma classifies → policy / financial / legal / entity / research",
                 "Below 0.5 confidence → keyword-scoring rule fallback",
                 "Images: 2-phase pipeline (sequential prefilter → parallel VLM)"
             ])

    add_card(slide, Inches(6.8), Inches(3.8), Inches(5.8), Inches(3.0),
             "Chunking, Embedding & Storing", [
                 "Legal → Gemma clause extraction with coverage-retry",
                 "Financial → semantic chunks, tables handled separately",
                 "Others → semantic breakpoint chunking (embedding-distance)",
                 "BGE bge-large-en-v1.5 embeddings, 1024 dimensions",
                 "Text → vector/clause/document store by type",
                 "Tables → VLM reconstruct + faithfulness gate → table_store"
             ])

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — Semantic Chunking (Deep Dive)
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "05", "Semantic Chunking — Deep Dive",
                          "Default path (CHUNK_USE_SEMANTIC=True) — splits at real topic shifts using sentence "
                          "embeddings, not fixed byte counts. Legal and legacy paths bypass this entirely.")

    chunk_steps = [
        ("1", "Group Units", "By section / list /\nimage boundaries"),
        ("2", "Embed Sentences", "BGE embeds each\nsentence in a unit"),
        ("3", "Find Breakpoints", "Split where cosine\ndistance > 95th pct"),
        ("4", "Force-Split", "Oversized (>1024 tok)\nat strongest boundary"),
        ("5", "Merge Small", "Undersized (<50 tok)\nmerges into neighbor"),
        ("6", "Gemma Enrich", "section_title, keywords,\nsemantic_type (batched)"),
    ]

    x = Inches(0.4)
    for num, title, desc in chunk_steps:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(1.95), Inches(1.15))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1.5)

        nbadge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.78), Inches(1.95), Inches(0.35), Inches(0.35))
        nbadge.fill.solid()
        nbadge.fill.fore_color.rgb = ACCENT
        nbadge.line.fill.background()
        ntf = nbadge.text_frame
        ntf.paragraphs[0].text = num
        ntf.paragraphs[0].font.size = Pt(10)
        ntf.paragraphs[0].font.color.rgb = TRUE_WHITE
        ntf.paragraphs[0].font.bold = True
        ntf.paragraphs[0].font.name = "Consolas"
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE

        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(9)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.name = "Calibri"

        if num != "6":
            add_text_box(slide, x + Inches(1.95), Inches(2.65), Inches(0.25), Inches(0.3),
                         "→", font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)
        x += Inches(2.05)

    add_card(slide, Inches(0.6), Inches(3.75), Inches(5.9), Inches(2.9),
             "No Overlap — Breadcrumb Instead", [
                 "No fixed-size overlap copied between consecutive chunks",
                 "Each chunk prefixed with a \"Context: A > B > C\" breadcrumb",
                 "Breadcrumb re-establishes hierarchical position for free",
                 "Cheaper than overlap; keeps boundaries clean at topic shifts",
                 "Chunk-level metadata carries section_title + keywords too"
             ])

    add_card(slide, Inches(6.75), Inches(3.75), Inches(5.9), Inches(2.9),
             "Per-Document-Type Paths", [
                 "Legal docs → clause-based chunking instead (bypasses this path)",
                 "Financial docs → same semantic path, minus image captions",
                 "Gemma enrichment runs in batches over all chunks",
                 "Falls back to heuristic tagging if Gemma call fails",
                 "Legacy fixed-512-token + overlap chunker: config flag opt-in"
             ], title_color=ORANGE)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 7 — Image OCR → Embedding Pipeline
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "06", "Image OCR → Embedding Pipeline",
                          "Every extracted image is prefiltered, OCR'd or VLM-analysed, embedded with the same "
                          "BGE model as text chunks, then stored with a signed URL for citation.")

    img_steps = [
        ("1", "Parse", "Extract images\nfrom PDF/docs"),
        ("2", "Prefilter", "Skip blank, tiny,\nduplicate images"),
        ("3", "Classify", "Plain text vs.\ntable/chart/complex"),
        ("4", "OCR / VLM", "OCR-only text, or\nVLM caption+structure"),
        ("5", "Build Text", "VLM text, fallback\nto raw OCR"),
        ("6", "Embed", "BGE bge-large-en-v1.5\n(same model as chunks)"),
        ("7", "Store", "image_store +\nsigned image URL"),
    ]

    x = Inches(0.25)
    for num, title, desc in img_steps:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.15), Inches(1.65), Inches(1.05))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT if num != "7" else GREEN
        box.line.width = Pt(1.5)

        nbadge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.63), Inches(1.9), Inches(0.35), Inches(0.35))
        nbadge.fill.solid()
        nbadge.fill.fore_color.rgb = ACCENT if num != "7" else GREEN
        nbadge.line.fill.background()
        ntf = nbadge.text_frame
        ntf.paragraphs[0].text = num
        ntf.paragraphs[0].font.size = Pt(10)
        ntf.paragraphs[0].font.color.rgb = TRUE_WHITE
        ntf.paragraphs[0].font.bold = True
        ntf.paragraphs[0].font.name = "Consolas"
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE

        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(3)
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(8)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.name = "Calibri"

        if num != "7":
            add_text_box(slide, x + Inches(1.65), Inches(2.55), Inches(0.2), Inches(0.3),
                         "→", font_size=13, color=ACCENT, alignment=PP_ALIGN.CENTER)
        x += Inches(1.72)

    add_card(slide, Inches(0.6), Inches(3.45), Inches(5.9), Inches(3.2),
             "Coverage — 46/46 Tests Passing (28 + 18)", [
                 "Prefiltering: blank / tiny-icon / duplicate detection (rule engine)",
                 "OCR/VLM analysis: text-only images → OCR; tables/charts → VLM",
                 "Embedding text build: VLM structured text preferred, OCR fallback",
                 "Embedding: BAAI/bge-large-en-v1.5 — identical model used for chunks",
                 "Storage: embedding vector + OCR/VLM text + page/doc_id metadata",
                 "  + signed URL back to the source image in Supabase Storage"
             ])

    add_card(slide, Inches(6.75), Inches(3.45), Inches(5.9), Inches(3.2),
             "Retrieval & Citation (Query Time)", [
                 "image_store is NOT searched directly — excluded from retriever",
                 "  store keys entirely (confirmed: test_retriever_images.py)",
                 "Visual queries route to whichever store holds the embedded",
                 "  OCR/caption text instead (vector / table / clause / document)",
                 "Citations resolve back to the source image's signed URL",
                 "Signed-URL generation failure is non-fatal to the citation"
             ], title_color=TEAL)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 8 — Multi-Store Routing
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "07", "Multi-Store Routing",
                          "Two independent routing decisions: which store a chunk is written to, and which store(s) a query searches.")

    # Routing table
    routes = [
        ("policy", "vector_store", "Semantic chunks, HNSW search"),
        ("financial", "table_store + vector_store", "Tables as JSON/markdown + surrounding text"),
        ("legal", "clause_store", "Clause-level, with risk / parties / dates"),
        ("research", "document_store", "Chunks with citation metadata"),
        ("entity", "vector_store", "Fallback when no specialised store fits"),
        ("tables (any)", "table_store / table_chunk_store", "Always routed here regardless of parent doc type"),
        ("images", "content-matched store", "Cross-stored + embedded by image_router"),
    ]

    # Table header
    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.6), Inches(1.9), Inches(12.0), Inches(0.4))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = ACCENT
    header_bg.line.fill.background()

    headers = [("Document Type", Inches(0.7), Inches(2.5)),
               ("Primary Store", Inches(3.3), Inches(3.5)),
               ("Notes", Inches(7.0), Inches(5.5))]
    for text, left, width in headers:
        add_text_box(slide, left, Inches(1.9), width, Inches(0.4),
                     text, font_size=11, color=TRUE_WHITE, bold=True)

    y = Inches(2.35)
    for i, (dtype, store, notes) in enumerate(routes):
        bg_color = CARD_BG if i % 2 == 0 else RGBColor(0xEB, 0xEC, 0xF3)
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.6), y, Inches(12.0), Inches(0.4))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = bg_color
        row_bg.line.fill.background()

        add_text_box(slide, Inches(0.7), y, Inches(2.5), Inches(0.4),
                     dtype, font_size=11, color=ORANGE, bold=True, font_name="Consolas")
        add_text_box(slide, Inches(3.3), y, Inches(3.5), Inches(0.4),
                     store, font_size=11, color=GREEN, font_name="Consolas")
        add_text_box(slide, Inches(7.0), y, Inches(5.5), Inches(0.4),
                     notes, font_size=11, color=LIGHT_GRAY)
        y += Inches(0.4)

    # Query-time store selection
    add_text_box(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(0.4),
                 "Query-Time Store Selection — intent_service.classify_intent",
                 font_size=16, color=WHITE, bold=True)

    intent_steps = [
        ("1  Rule-based keyword hit", "Wins outright if confidence ≥ 0.8"),
        ("2  Semantic router", "Cosine similarity vs. 4 per-store prototype centroids"),
        ("3  Gemma LLM call", "Off critical path, gated by INTENT_USE_LLM"),
    ]
    x = Inches(0.6)
    for step, desc in intent_steps:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.9), Inches(3.8), Inches(0.75))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = step
        run.font.size = Pt(12)
        run.font.color.rgb = ACCENT_LIGHT
        run.font.bold = True
        run.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(9)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.name = "Calibri"
        x += Inches(4.1)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 9 — Query Routing (Deep Dive)
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "08", "Query Routing — Deep Dive",
                          "Which store(s) get searched — decided before any embedding search runs. Stops at the "
                          "first confident tier; falls back to searching everything when nothing is sure.")

    route_tiers = [
        ("0", "Explicit Type", "document_types passed\nin → wins outright", ACCENT_LIGHT),
        ("1", "Rules", "Keyword match, e.g.\n\"clause\"→clause,\n\"revenue\"→table+vector", ACCENT),
        ("2", "Semantic Router", "Cosine sim vs. per-store\nprototype centroids —\ncatches paraphrases", TEAL),
        ("3", "Gemma LLM", "Off by default — full\nclassification call,\nonly if 1–2 unsure", ORANGE),
        ("4", "Fallback", "Confidence < 0.5 →\nsearch all 4 stores\n(recall-safe)", RED_SOFT),
    ]

    x = Inches(0.4)
    for num, title, desc, color in route_tiers:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(2.3), Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(1.5)

        nbadge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.95), Inches(1.95), Inches(0.35), Inches(0.35))
        nbadge.fill.solid()
        nbadge.fill.fore_color.rgb = color
        nbadge.line.fill.background()
        ntf = nbadge.text_frame
        ntf.paragraphs[0].text = num
        ntf.paragraphs[0].font.size = Pt(10)
        ntf.paragraphs[0].font.color.rgb = TRUE_WHITE
        ntf.paragraphs[0].font.bold = True
        ntf.paragraphs[0].font.name = "Consolas"
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE

        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(9)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.name = "Calibri"

        if num != "4":
            add_text_box(slide, x + Inches(2.3), Inches(2.85), Inches(0.25), Inches(0.3),
                         "→", font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)
        x += Inches(2.42)

    add_card(slide, Inches(0.6), Inches(4.15), Inches(12.05), Inches(2.5),
             "Concurrent Multi-Store Search (Once Stores Are Selected)", [
                 "Every selected store runs concurrently — thread pool, one dedicated DB connection each",
                 "Each store independently does its own HNSW ANN search (pgvector, cosine distance)",
                 "All per-store result lists merge and sort by distance into one candidate pool",
                 "Candidate pool then proceeds to cross-encoder reranking (bge-reranker-large)",
                 "Same tiered logic (rules → semantic router → Gemma) decides write-time routing too —",
                 "  see the Multi-Store Routing table for the document-type → store mapping"
             ])

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 10 — Query & Retrieval Pipeline
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "09", "Query & Retrieval Pipeline",
                          "POST /api/v1/query and the streaming variant /query/stream — full RAG query flow.")

    query_steps = [
        ("1", "Conversational\nShort-Circuit", "Greeting / small-talk\nskips retrieval", MEDIUM_GRAY),
        ("2", "Retrieval\n(Classic or Agentic)", "Hybrid search or\nRAVEN/SPYDER loop", ACCENT),
        ("3", "GraphRAG\nFusion", "Local / global / none\nmulti-hop expansion", TEAL),
        ("4", "Structured\nTable Query", "SUM/AVG/COUNT/MIN/MAX\nregex over json_data", ORANGE),
        ("5", "Ranking", "Cross-encoder rerank\nRRF fusion", ACCENT_LIGHT),
        ("6", "Synthesis", "Gemma JSON answer\nconfidence blending", GREEN),
    ]

    x = Inches(0.3)
    for num, title, desc, color in query_steps:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(2.0), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Number
        nbadge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.75), Inches(1.95), Inches(0.4), Inches(0.4))
        nbadge.fill.solid()
        nbadge.fill.fore_color.rgb = color
        nbadge.line.fill.background()
        ntf = nbadge.text_frame
        ntf.paragraphs[0].text = num
        ntf.paragraphs[0].font.size = Pt(12)
        ntf.paragraphs[0].font.color.rgb = TRUE_WHITE
        ntf.paragraphs[0].font.bold = True
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE

        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(9)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.name = "Calibri"

        if int(num) < 6:
            add_text_box(slide, x + Inches(2.0), Inches(2.9), Inches(0.2), Inches(0.3),
                         "→", font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)
        x += Inches(2.15)

    # Hybrid search detail
    add_card(slide, Inches(0.6), Inches(4.5), Inches(6.0), Inches(2.4),
             "Hybrid Search = Semantic + Keyword, Fused", [
                 "Semantic: pgvector HNSW cosine search per selected store",
                 "Keyword: Postgres FTS (ts_rank_cd / websearch_to_tsquery)",
                 "Fusion: Reciprocal Rank Fusion (k=60) merges ranked lists",
                 "Balanced pool: top-8 per store, graph hits exempted from cap",
                 "BGE cross-encoder reranker with sigmoid-blended scores"
             ])

    # Synthesis detail
    add_card(slide, Inches(6.95), Inches(4.5), Inches(5.8), Inches(2.4),
             "Synthesis — Gemma 4 Answer Generation", [
                 "Numbered, per-store context blocks → markdown tables",
                 "Images → caption + OCR text, budget-capped at 12K chars",
                 "Strict JSON prompt: {answer, confidence, sources_used, notes}",
                 "Confidence = Gemma self-rating (40%) + retrieval signal (60%)",
                 "Streaming route: live tokens + post-hoc confidence rating"
             ])

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 11 — GraphRAG Layer
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "10", "GraphRAG Layer",
                          "A Neo4j knowledge graph layered on top of Postgres stores — entity extraction, "
                          "relationship mapping, and community detection for multi-hop retrieval.")

    # Graph model visualization
    nodes = [
        ("Document", "id, filename, doc_type", Inches(0.8), ORANGE),
        ("Entity", "key (dedup), name, type", Inches(4.6), ACCENT_LIGHT),
        ("Entity", "(target node)", Inches(8.0), ACCENT_LIGHT),
        ("Community", "level, title, summary", Inches(10.8), GREEN),
    ]

    for name, props, x, color in nodes:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(2.0), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = name
        run.font.size = Pt(14)
        run.font.color.rgb = color
        run.font.bold = True
        run.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = props
        run2.font.size = Pt(9)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.name = "Consolas"

    # Edge labels
    edges = [
        ("← MENTIONED_IN →", Inches(3.0)),
        ("—[ relationship ]→", Inches(6.8)),
        ("IN_COMMUNITY →", Inches(10.2)),
    ]
    for text, x in edges:
        add_text_box(slide, x, Inches(2.6), Inches(1.8), Inches(0.4),
                     text, font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER,
                     font_name="Consolas")

    # Local vs Global search
    add_card(slide, Inches(0.6), Inches(3.8), Inches(5.8), Inches(3.0),
             "🔍  Local Search", [
                 "Triggered when query mentions specific entities",
                 "Extracts entities from query text",
                 "Walks graph up to 2 hops from matched entities",
                 "Hydrates pg_id references back to full Postgres chunks",
                 "Finds cross-document connections via shared entities",
                 "Example: Drug X → CONTAINS → Penicillin → MAY_CAUSE → Allergy"
             ], title_color=ACCENT_LIGHT)

    add_card(slide, Inches(6.8), Inches(3.8), Inches(5.8), Inches(3.0),
             "🌐  Global Search (MAP-REDUCE)", [
                 "Triggered by aggregation questions without specific entities",
                 "Example: 'What are the main themes across all contracts?'",
                 "MAP: Gemma rates relevance + drafts partial answers per community",
                 "REDUCE: Gemma synthesises final answer over all partial answers",
                 "Community summaries optionally embedded for vector similarity",
                 "Louvain clustering with label-propagation fallback"
             ], title_color=GREEN)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 12 — GraphRAG: Build Time & Query Time (Deep Dive)
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "11", "GraphRAG — Build Time & Query Time",
                          "run_graph_stage() at ingestion writes the graph; route_graphrag() at query time decides "
                          "local-hop, global map-reduce, or a lighter cross-document expansion.")

    add_card(slide, Inches(0.5), Inches(2.0), Inches(6.1), Inches(4.9),
             "Build Time — Ingestion Stage 6 (run_graph_stage)", [
                 "Clear existing graph for this doc (re-ingest safety), upsert Document node",
                 "assemble_chunk_records() bridges Postgres UUIDs (vector/clause/",
                 "  document_store) to their chunk text",
                 "Per-chunk extraction: parallel ThreadPoolExecutor, capped at",
                 "  GRAPHRAG_EXTRACT_CONCURRENCY=4 — one Gemma 4 call per chunk pulls",
                 "  entities (org/person/location/product/law/project/concept) + typed",
                 "  relationships as JSON",
                 ">200 chunks (GRAPHRAG_MAX_CHUNKS_PER_DOC) → falls back to a single",
                 "  doc-level extraction instead of per-chunk",
                 "Entity names canonicalized + MERGEd (dedup across chunks/documents);",
                 "  relationships written between canonical entity keys",
                 "Writes are verified — counts reflect chunks that actually landed in",
                 "  Neo4j, not extraction attempts",
                 "Additive extraction also runs over table_store / image_store rows",
                 "  (markdown / OCR / VLM text) so tables and images contribute entities",
                 "Redis \"dirty\" counter bumped → recompute_communities_task enqueued",
                 "  (debounced: fires after GRAPHRAG_COMMUNITY_MIN_INTERVAL_SEC=1800s OR",
                 "  GRAPHRAG_COMMUNITY_DIRTY_DOCS=5 dirty docs) → Louvain community",
                 "  detection + per-community summaries/embeddings for global search"
             ], title_color=ACCENT_LIGHT)

    add_card(slide, Inches(6.75), Inches(2.0), Inches(6.1), Inches(4.9),
             "Query Time — route_graphrag(query)", [
                 "Extracts entities from the query, checks if any exist as nodes in",
                 "  Neo4j (fuzzy substring match) → mode = \"local\"",
                 "Else: query has aggregation cue words (\"summarize\", \"compare\",",
                 "  \"overall\", \"trends\"...) → mode = \"global\"; else → mode = \"none\"",
                 "",
                 "LOCAL: canonicalize query entities → local_neighborhood() walks",
                 "  GRAPHRAG_LOCAL_HOPS=2 hops across relationship edges, then back out",
                 "  via MENTIONED_IN to get {pg_id, document_id, store} — these extra",
                 "  chunks are deduped and merged into the normal vector-search",
                 "  rerank pool",
                 "",
                 "GLOBAL: select top communities (by query-embedding similarity to",
                 "  summaries, or top-N) → MAP: Gemma per-community \"is this relevant,",
                 "  partial answer?\" concurrently → REDUCE: Gemma synthesises final",
                 "  answer from partial answers. Bypasses vector retrieval/reranking",
                 "  entirely and returns directly",
                 "",
                 "NONE: a lighter multi-PDF graph expansion still runs on retrieved",
                 "  chunks (unless a single document_id is pinned), pulling in related",
                 "  chunks from other documents that share entities"
             ], title_color=GREEN)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 13 — Agentic RAG Loop
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "12", "Agentic RAG Loop",
                          "RAVEN/SPYDER loop — iterative query refinement with sufficiency judgment. "
                          "Gated by AGENTIC_RAG_ENABLED, replaces classic retrieval when enabled.")

    # Agentic flow
    agentic_steps = [
        ("GraphRAG\nGlobal Check", TEAL),
        ("RAVEN\nReframes Query", ACCENT),
        ("Hybrid\nRetrieve", CARD_BG),
        ("Graph\nExpansion", TEAL),
        ("Rerank", CARD_BG),
        ("SPYDER\nJudges", ORANGE),
    ]

    x = Inches(0.5)
    for title, color in agentic_steps:
        fill_color = color if color != CARD_BG else CARD_BG
        border_color = color if color != CARD_BG else ACCENT
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(1.8), Inches(1.0))
        box.fill.solid()
        if color == CARD_BG:
            box.fill.fore_color.rgb = CARD_BG
        else:
            box.fill.fore_color.rgb = RGBColor(
                max(0, color[0] - 0x50),
                max(0, color[1] - 0x50),
                max(0, color[2] - 0x50)
            )
        box.line.color.rgb = border_color
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Calibri"

        if x < Inches(10):
            add_text_box(slide, x + Inches(1.8), Inches(2.6), Inches(0.3), Inches(0.3),
                         "→", font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)
        x += Inches(2.1)

    # Loop back arrow indicator
    add_text_box(slide, Inches(3.0), Inches(3.5), Inches(8.0), Inches(0.4),
                 "↺  Not sufficient & loop budget remains → reframed query loops back (max 2 iterations)",
                 font_size=12, color=ORANGE, alignment=PP_ALIGN.CENTER)

    # RAVEN and SPYDER cards
    add_card(slide, Inches(0.6), Inches(4.2), Inches(5.8), Inches(2.5),
             "🦅  RAVEN — Pre-Retrieval Agent", [
                 "Gemma call that reframes the raw query into cleaner search form",
                 "Optionally decomposes into up to 3 sub-queries",
                 "Emits a store hint for targeted retrieval",
                 "Fails open to raw query if disabled or on parse failure",
                 "Reduces ambiguity before vector/keyword search"
             ], title_color=ACCENT_LIGHT)

    add_card(slide, Inches(6.8), Inches(4.2), Inches(5.8), Inches(2.5),
             "🕷  SPYDER — Post-Rerank Judge", [
                 "Given query + numbered reranked context, decides sufficiency",
                 "Returns: confidence score, what's missing, reframed query",
                 "Fails open to 'sufficient' on any error — loop can never hang",
                 "Stop: sufficient, confidence ≥ 0.6, no reframe, or cap reached",
                 "Ensures high-quality answers through iterative refinement"
             ], title_color=ORANGE)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 14 — API Reference
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "13", "API Reference",
                          "All routes under /api/v1 — Documents, Ingestion, Query, Graph, Chats & Health endpoints.")

    api_groups = [
        ("Documents", [
            ("GET", "/documents", "Paginated list, filter by status/type"),
            ("GET", "/documents/{id}", "Full detail + per-store chunk counts"),
            ("GET", "/documents/{id}/page-stats", "Per-page chunk/table/clause/image counts"),
            ("GET", "/documents/{id}/images", "Extracted images + signed URLs"),
            ("GET", "/documents/{id}/chunks", "Paginated chunks from specific store"),
            ("DELETE", "/documents/{id}", "Cascades: DB + Storage + Neo4j cleanup"),
        ]),
        ("Ingest", [
            ("POST", "/ingest/upload", "Single-file upload → dispatch pipeline"),
            ("POST", "/ingest/pipeline", "Batch upload with taxonomy metadata"),
            ("GET", "/ingest/pipelines", "List pipeline runs"),
            ("POST", "/ingest/documents/{id}/reprocess", "Re-dispatch ingestion"),
            ("GET", "/ingest/status/{job_id}", "Poll ingestion job progress"),
        ]),
    ]

    x_start = Inches(0.6)
    for col_idx, (group_name, endpoints) in enumerate(api_groups):
        x = x_start + col_idx * Inches(6.3)

        # Group title
        add_text_box(slide, x, Inches(1.9), Inches(5.8), Inches(0.4),
                     group_name, font_size=14, color=ACCENT_LIGHT, bold=True)

        y = Inches(2.35)
        for method, path, desc in endpoints:
            # Method badge color
            method_colors = {"GET": GREEN, "POST": ACCENT_LIGHT, "DELETE": RED_SOFT, "PATCH": ORANGE}
            mc = method_colors.get(method, LIGHT_GRAY)

            # Method badge
            mbadge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            x, y, Inches(0.6), Inches(0.28))
            mbadge.fill.solid()
            mbadge.fill.fore_color.rgb = RGBColor(
                min(255, mc[0] + (255 - mc[0]) * 3 // 4), min(255, mc[1] + (255 - mc[1]) * 3 // 4), min(255, mc[2] + (255 - mc[2]) * 3 // 4)
            )
            mbadge.line.fill.background()
            mtf = mbadge.text_frame
            mtf.paragraphs[0].text = method
            mtf.paragraphs[0].font.size = Pt(8)
            mtf.paragraphs[0].font.color.rgb = mc
            mtf.paragraphs[0].font.bold = True
            mtf.paragraphs[0].font.name = "Consolas"
            mtf.paragraphs[0].alignment = PP_ALIGN.CENTER
            mtf.vertical_anchor = MSO_ANCHOR.MIDDLE

            add_text_box(slide, x + Inches(0.7), y, Inches(2.5), Inches(0.28),
                         path, font_size=10, color=WHITE, font_name="Consolas")
            add_text_box(slide, x + Inches(3.2), y, Inches(2.5), Inches(0.28),
                         desc, font_size=9, color=LIGHT_GRAY)
            y += Inches(0.35)

    # Additional API groups below
    more_groups = [
        ("Query", [
            ("POST", "/query", "Full RAG query → JSON answer + citations"),
            ("POST", "/query/stream", "SSE streaming — token deltas + stage events"),
        ]),
        ("Graph", [
            ("GET", "/graph/status", "Neo4j connectivity + node/edge counts"),
            ("GET", "/graph/entities", "Paginated Entity node list"),
            ("GET", "/graph/search", "Multi-hop entity search"),
            ("POST", "/graph/rebuild/{id}", "Rebuild document's graph contribution"),
        ]),
    ]

    for col_idx, (group_name, endpoints) in enumerate(more_groups):
        x = x_start + col_idx * Inches(6.3)
        y_base = Inches(4.7)
        add_text_box(slide, x, y_base, Inches(5.8), Inches(0.4),
                     group_name, font_size=14, color=ACCENT_LIGHT, bold=True)

        y = y_base + Inches(0.45)
        for method, path, desc in endpoints:
            method_colors = {"GET": GREEN, "POST": ACCENT_LIGHT, "DELETE": RED_SOFT, "PATCH": ORANGE}
            mc = method_colors.get(method, LIGHT_GRAY)

            mbadge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            x, y, Inches(0.6), Inches(0.28))
            mbadge.fill.solid()
            mbadge.fill.fore_color.rgb = RGBColor(
                min(255, mc[0] // 4), min(255, mc[1] // 4), min(255, mc[2] // 4)
            )
            mbadge.line.fill.background()
            mtf = mbadge.text_frame
            mtf.paragraphs[0].text = method
            mtf.paragraphs[0].font.size = Pt(8)
            mtf.paragraphs[0].font.color.rgb = mc
            mtf.paragraphs[0].font.bold = True
            mtf.paragraphs[0].font.name = "Consolas"
            mtf.paragraphs[0].alignment = PP_ALIGN.CENTER
            mtf.vertical_anchor = MSO_ANCHOR.MIDDLE

            add_text_box(slide, x + Inches(0.7), y, Inches(2.5), Inches(0.28),
                         path, font_size=10, color=WHITE, font_name="Consolas")
            add_text_box(slide, x + Inches(3.2), y, Inches(2.5), Inches(0.28),
                         desc, font_size=9, color=LIGHT_GRAY)
            y += Inches(0.35)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 15 — Database Schema
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "14", "Database Schema",
                          "Postgres schema multi_store_rag_working — 16 migrations. Plus a separate Neo4j graph store.")

    tables = [
        ("document_registry", "Central registry — one row per uploaded document", ACCENT_LIGHT),
        ("vector_store", "Policy / generic semantic chunks with 1024-dim embeddings", ACCENT_LIGHT),
        ("table_store", "All extracted tables — JSON, CSV, markdown formats", ORANGE),
        ("table_chunk_store", "Row-window children of large tables", ORANGE),
        ("clause_store", "Legal clauses with risk_level, obligor/obligee, key_dates", RED_SOFT),
        ("document_store", "Research chunks with citation_key, authors, doi, journal", GREEN),
        ("image_store", "Image repository — no embedding column, cross-embedded", TEAL),
        ("ingestion_jobs", "Live pipeline progress, polled by frontend", MEDIUM_GRAY),
        ("pipeline_runs", "Batch upload grouping with taxonomy", MEDIUM_GRAY),
        ("chat_sessions / messages", "Chat history with citations and confidence", MEDIUM_GRAY),
    ]

    # Layout: 2 columns of cards
    for i, (table_name, desc, color) in enumerate(tables):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + col * Inches(6.3)
        y = Inches(2.0) + row * Inches(0.85)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(1.5)

        add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.5), Inches(0.3),
                     table_name, font_size=12, color=color, bold=True, font_name="Consolas")
        add_text_box(slide, x + Inches(0.15), y + Inches(0.35), Inches(5.5), Inches(0.3),
                     desc, font_size=10, color=LIGHT_GRAY)

    # Connection pooling note
    add_text_box(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
                 "Connection pooling: psycopg2.pool.ThreadedConnectionPool (min 2, max 20), pgvector registered per-connection, semaphore-bounded with 8s timeout",
                 font_size=10, color=MEDIUM_GRAY)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 16 — Config & Feature Flags
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "15", "Config & Feature Flags",
                          "backend/app/config.py declares every advanced capability OFF by default. "
                          "Feature flags control graceful degradation.")

    flags_on = [
        "AGENTIC_RAG_ENABLED", "GRAPHRAG_ENABLED", "NEO4J_ENABLED",
        "INGESTION_STAGED_ENABLED", "HYBRID_IN_CLASSIC_PATH", "INTENT_USE_LLM",
        "RETRIEVAL_CACHE_ENABLED", "RAVEN_ENABLED", "SPYDER_ENABLED",
        "HYBRID_SEARCH_ENABLED", "STRUCTURED_QUERY_ENABLED", "TABLE_CHILD_SEARCH_ENABLED",
        "CHUNK_USE_SEMANTIC", "INTENT_USE_SEMANTIC_ROUTER", "TABLE_VLM_RECONSTRUCT",
        "PREFILTER_ENABLED"
    ]

    flags_off = ["RETRIEVAL_CACHE_RESULTS_ENABLED"]

    # ON flags grid
    x = Inches(0.6)
    y = Inches(2.0)
    for i, flag in enumerate(flags_on):
        col = i % 4
        row = i // 4
        fx = x + col * Inches(3.1)
        fy = y + row * Inches(0.55)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, fx, fy, Inches(2.85), Inches(0.42))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = RGBColor(0xD8, 0xDA, 0xE5)
        box.line.width = Pt(1)

        add_text_box(slide, fx + Inches(0.1), fy + Inches(0.05), Inches(2.0), Inches(0.32),
                     flag, font_size=8, color=WHITE, font_name="Consolas")

        # ON pill
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      fx + Inches(2.2), fy + Inches(0.08), Inches(0.5), Inches(0.25))
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(0xE4, 0xF7, 0xEE)
        pill.line.fill.background()
        ptf = pill.text_frame
        ptf.paragraphs[0].text = "ON"
        ptf.paragraphs[0].font.size = Pt(8)
        ptf.paragraphs[0].font.color.rgb = GREEN
        ptf.paragraphs[0].font.bold = True
        ptf.paragraphs[0].font.name = "Consolas"
        ptf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # OFF flag
    for flag in flags_off:
        fy = y + Inches(2.3)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, fy, Inches(3.5), Inches(0.42))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = RGBColor(0xF5, 0xD5, 0xD0)
        box.line.width = Pt(1)

        add_text_box(slide, x + Inches(0.1), fy + Inches(0.05), Inches(2.6), Inches(0.32),
                     flag, font_size=8, color=WHITE, font_name="Consolas")

        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x + Inches(2.8), fy + Inches(0.08), Inches(0.5), Inches(0.25))
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(0xFB, 0xE9, 0xE7)
        pill.line.fill.background()
        ptf = pill.text_frame
        ptf.paragraphs[0].text = "OFF"
        ptf.paragraphs[0].font.size = Pt(8)
        ptf.paragraphs[0].font.color.rgb = RED_SOFT
        ptf.paragraphs[0].font.bold = True
        ptf.paragraphs[0].font.name = "Consolas"
        ptf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Settings groups
    settings_groups = [
        ("Gemma 4 / CDAC", ["Model: gemma-4-27b-it", "Max tokens: 800", "Max concurrent: 5", "Retries + backoff on 429/5xx"]),
        ("Embeddings / Rerank", ["BAAI/bge-large-en-v1.5", "ms-marco-MiniLM-L-6-v2", "Batch size 32, CPU device"]),
        ("Multi-hop / GraphRAG", ["Max hops: 2", "Entities/chunk: 15", "Community: Louvain", "Recompute interval: 30 min"]),
        ("Docling / VLM", ["OCR + table structure: on", "Image scale: 2.0×", "VLM concurrency: 4", "Context cap: 12K chars"]),
    ]

    x = Inches(0.6)
    for group_title, items in settings_groups:
        add_card(slide, x, Inches(5.0), Inches(2.85), Inches(2.0),
                 group_title, items, title_color=ACCENT_LIGHT)
        x += Inches(3.1)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 17 — Frontend
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_section_header(slide, "16", "Frontend",
                          "Next.js 14 App Router — full-featured UI for document management, ingestion monitoring, and chat-based queries.")

    # Pages
    pages = [
        ("/", "Landing / Home", "Entry point"),
        ("/upload", "Upload UI", "Drag-drop + live per-stage progress"),
        ("/documents", "Document Browser", "Filterable list + detail view"),
        ("/documents/[id]", "Document Detail", "Chunks, images, page stats"),
        ("/pipelines", "Pipeline History", "Batch run tracking"),
        ("/query", "Chat / Query", "Full conversational RAG interface"),
    ]

    add_text_box(slide, Inches(0.6), Inches(2.0), Inches(3), Inches(0.4),
                 "Pages", font_size=16, color=WHITE, bold=True)

    y = Inches(2.5)
    for route, name, desc in pages:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.6), y, Inches(5.5), Inches(0.55))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = RGBColor(0xD8, 0xDA, 0xE5)
        box.line.width = Pt(1)

        add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(1.5), Inches(0.25),
                     route, font_size=10, color=ACCENT_LIGHT, font_name="Consolas")
        add_text_box(slide, Inches(2.4), y + Inches(0.05), Inches(1.8), Inches(0.25),
                     name, font_size=11, color=WHITE, bold=True)
        add_text_box(slide, Inches(0.8), y + Inches(0.28), Inches(5.0), Inches(0.25),
                     desc, font_size=9, color=LIGHT_GRAY)
        y += Inches(0.62)

    # Key components
    add_card(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.7),
             "Key Components & Libraries", [
                 "ChunkViewer, DocumentTable — data display",
                 "ChunkingDetail, ImagesDetail, ParsingDetail — per-stage panels",
                 "FileDropzone, UploadProgress — upload flow",
                 "lib/api-client.ts — typed client wrapping every endpoint",
                 "hooks/useJobStatus.ts — polling hook for ingestion progress",
                 "",
                 "TanStack Query for server-state / polling",
                 "Tailwind CSS for styling",
                 "API proxy routes forward to FastAPI backend",
                 "",
                 "Note: No graph-view UI yet — GraphRAG surface is",
                 "currently backend/API-only"
             ], title_color=TEAL)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 18 — Summary / Thank You
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Decorative top bar
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    deco.fill.solid()
    deco.fill.fore_color.rgb = ACCENT
    deco.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
                 "DECISION MINDS", font_size=14, color=ACCENT_LIGHT, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(1.0),
                 "Multi-Store RAG Chatbot", font_size=44, color=DARK_BG, bold=True)

    add_text_box(slide, Inches(0.8), Inches(2.5), Inches(10), Inches(0.5),
                 "System Architecture Summary", font_size=20, color=LIGHT_GRAY)

    # Key highlights
    highlights = [
        ("5 Specialised Stores", "Type-aware routing maximises retrieval precision"),
        ("GraphRAG + Neo4j", "Multi-hop entity connections across documents"),
        ("Agentic RAVEN/SPYDER", "Iterative query refinement with sufficiency judgment"),
        ("Hybrid Search", "Semantic + keyword fusion with cross-encoder reranking"),
        ("Graceful Degradation", "Every smart stage has a rule-based fallback"),
        ("Full Observability", "Live pipeline progress, per-stage timing, confidence scores"),
    ]

    y = Inches(3.3)
    for i, (title, desc) in enumerate(highlights):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.0)
        yy = y + row * Inches(0.75)

        # Accent dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, yy + Inches(0.05), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()

        add_text_box(slide, x + Inches(0.3), yy, Inches(2.5), Inches(0.3),
                     title, font_size=13, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.3), yy + Inches(0.25), Inches(5.0), Inches(0.3),
                     desc, font_size=10, color=LIGHT_GRAY)

    add_text_box(slide, Inches(0.8), Inches(6.3), Inches(10), Inches(0.4),
                 "Thank You", font_size=24, color=ACCENT_LIGHT, bold=True)

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(5), Inches(0.4),
                 "Enterprise Multi-Store RAG System  ·  Decision Minds  ·  July 2026",
                 font_size=11, color=MEDIUM_GRAY)

    # ── Save ───────────────────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Multi_Store_RAG_Chatbot_Architecture.pptx")
    prs.save(output_path)
    print(f"\n[OK] Presentation saved to: {output_path}")
    print(f"     Slides: {len(prs.slides)}")
    print(f"     Format: Widescreen 16:9")
    return output_path


if __name__ == "__main__":
    build_presentation()
